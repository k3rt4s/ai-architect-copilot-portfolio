#!/usr/bin/env python3
# orchestrator.py
#
# PURPOSE
#   Ingest architecture documents (PDF, DOCX, PPTX, Excel, images, CSV, TXT, MD)
#   and run a 3-stage Gemini analysis pipeline:
#     Stage 1 — Corpus intake and structured indexing
#     Stage 2 — Architecture gap analysis
#     Stage 3 — Internet-enabled interactive chat advisor
#
# OUTPUTS
#   <out>/01_Index_<timestamp>.md
#   <out>/02_Gap_Analysis_<timestamp>.md
#   <out>/chat_transcript_<timestamp>.md
#
# RUN
#   python orchestrator.py --folder "<docs_dir>" --out "<output_dir>"

import sys
import os
import argparse
from pathlib import Path
from datetime import datetime

import pandas as pd
import docx
from pptx import Presentation
import pymupdf4llm
import shutil as _shutil

from dotenv import load_dotenv

load_dotenv()

_PACKAGE_DIR = Path(__file__).resolve().parent


def _resolve_tesseract() -> str:
    cmd = os.getenv("TESSERACT_CMD")
    if cmd:
        return cmd
    found = _shutil.which("tesseract")
    if found:
        return found
    raise EnvironmentError(
        "Tesseract not found. Set TESSERACT_CMD in your .env file, "
        "e.g. TESSERACT_CMD=C:/Program Files/Tesseract-OCR/tesseract.exe, "
        "or ensure tesseract is on your system PATH."
    )


import pytesseract
pytesseract.pytesseract.tesseract_cmd = _resolve_tesseract()

from PIL import Image

from gemini_runner import GeminiRunner


def extract_text(p: Path) -> str:
    ext = p.suffix.lower()
    try:
        content = ""
        if ext == ".pdf":
            content = pymupdf4llm.to_markdown(str(p))
        elif ext in {".xlsx", ".xls"}:
            df_dict = pd.read_excel(p, sheet_name=None)
            content = "\n".join(
                [f"### Sheet: {n}\n{df.to_markdown(index=False)}" for n, df in df_dict.items()]
            )
        elif ext == ".docx":
            doc = docx.Document(p)
            content = "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(p)
            content = "\n".join(
                [
                    shape.text
                    for slide in prs.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                ]
            )
        elif ext in {".png", ".jpg", ".jpeg"}:
            content = f"[OCR DATA]: {pytesseract.image_to_string(Image.open(p))}"
        elif ext in {".csv", ".txt", ".md"}:
            if ext == ".csv":
                content = pd.read_csv(p).to_markdown(index=False)
            else:
                content = p.read_text(encoding="utf-8", errors="ignore")

        return f"\n\n=== START OF FILE: {p.name} ===\n{content}\n=== END OF FILE: {p.name} ===\n"
    except Exception as e:
        return f"[Error extracting {p.name}: {str(e)}]"


def run_sequential_analysis(
    docs_text: str,
    file_list: list,
    project: str,
    location: str,
    model: str,
    output_dir: str,
) -> None:
    """Run the three-stage Gemini pipeline against the ingested corpus.

    Production behavior:
        Stage 1: load Persona_01 prompt, instantiate a GeminiRunner, feed it
                 the ingested document text, write the structured index to
                 ``01_Index_<timestamp>.md``.
        Stage 2: load Persona_02 prompt, feed the Stage 1 index, write the
                 gap-analysis report to ``02_Gap_Analysis_<timestamp>.md``.
        Stage 3: load Persona_03 prompt with the Google Search tool attached,
                 run an interactive chat loop against the gap analysis,
                 stream the transcript to ``chat_transcript_<timestamp>.md``.

    This portfolio sample stubs the stage orchestration intentionally. The
    surrounding architecture, the document-ingestion layer (`extract_text`),
    and the SDK wrapper (`gemini_runner.GeminiRunner`) are shown so the
    review-time pattern is auditable, but the tuned three-stage sequence and
    its persona prompts are proprietary.
    """
    raise NotImplementedError(
        "Stage orchestration withheld. This repository is a portfolio "
        "sample. The production three-stage Gemini pipeline and its persona "
        "prompts are proprietary."
    )


def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "Architecture document review and gap analysis. "
            "Portfolio sample: the orchestration stage is stubbed; see "
            "README.md and PORTFOLIO_SAMPLE.md for what is shown and what "
            "is withheld."
        )
    )
    parser.add_argument("--folder", required=True, help="Directory containing architecture documents")
    parser.add_argument("--out", required=True, help="Output directory for analysis reports")
    args = parser.parse_args()

    print(
        "ai-architect-copilot-portfolio: this repository is a portfolio "
        "sample. The three-stage Gemini pipeline is intentionally not "
        "included here. The CLI will ingest the folder (demonstrating the "
        "document-extraction layer) and then exit cleanly without "
        "performing analysis.",
        file=sys.stderr,
    )

    project = os.getenv("GOOGLE_CLOUD_PROJECT", os.getenv("VERTEX_PROJECT"))
    location = os.getenv("GOOGLE_CLOUD_LOCATION", os.getenv("VERTEX_LOCATION", "us-central1"))
    model = os.getenv("MODEL_NAME", "gemini-2.5-pro")

    folder = Path(args.folder).expanduser().resolve()
    if not folder.is_dir():
        print(f"X Folder not found: {folder}", file=sys.stderr)
        raise SystemExit(1)

    full_text = ""
    f_names: list = []

    for f in sorted(folder.rglob("*")):
        if f.is_file() and not f.name.startswith((".", "~$")):
            print(f"  → [INGESTING] {f.name}")
            data = extract_text(f)
            if data:
                full_text += data
                f_names.append(f.name)

    print(
        f"Ingested {len(f_names)} files; "
        f"{len(full_text):,} characters of extracted text. "
        "Stopping before the stubbed orchestration step. Exit code 2.",
        file=sys.stderr,
    )
    raise SystemExit(2)


if __name__ == "__main__":
    main()
